from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title, subtitle_text=None, bullet_points=None):
        slide_layout = prs.slide_layouts[1] if bullet_points else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle_text:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
        
        if bullet_points:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = bullet_points[0]
            for point in bullet_points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Student Assessment & Performance Tracker"
    subtitle.text = "Team Presentation - January 31, 2026\nModernizing academic performance tracking through automation and analytics."

    # 2. Problem Statement
    add_slide("Problem Statement", bullet_points=[
        "Manual Overhead: Faculty spend excessive time recording and calculating grades manually.",
        "Communication Gap: Parents often lack real-time updates on their child's progress.",
        "Data Fragmentation: Hard to track long-term performance trends across multiple assessments.",
        "Error Proneness: High risk of errors in manual mark entry and trend analysis."
    ])

    # 3. Proposed Solution
    add_slide("Proposed Solution", bullet_points=[
        "Full-Stack Platform: A unified portal for Admins, Faculty, Students, and Parents.",
        "Automation First: Automated grading, trend analysis, and multi-channel notifications.",
        "Scalability: Bulk upload functionality for large classrooms.",
        "Modern UI: Clean, intuitive interface for all user roles."
    ])

    # 4. Technical Architecture
    add_slide("Technical Architecture", bullet_points=[
        "Frontend: React (SPA) for high reactivity and modern UI.",
        "Backend: Spring Boot for robust business logic and secure API management.",
        "Database: MySQL for structured data storage.",
        "Integrations: Twilio (SMS) and SMTP (Email) for automated communication."
    ])

    # 5. Key Modules
    add_slide("Key Modules", bullet_points=[
        "Admin: The management hub (Users, Subjects, Departments).",
        "Faculty: The data engine (Marks Entry, Bulk Upload, Grade Calc).",
        "Student Dashboard: The insights portal (Scorecards, Progress Trends)."
    ])

    # 6. Unique Selling Points (USP)
    add_slide("Unique Selling Points (USP)", bullet_points=[
        "Automatic Trend Analysis: Comparing performance across time, not just individual scores.",
        "Hybrid Entry: Flexibility of manual entry or high-speed bulk CSV uploads.",
        "Multi-Stakeholder Alerts: Keeping both students and parents in the loop instantly."
    ])

    # 7. Results & Impact
    add_slide("Results & Impact", bullet_points=[
        "80% Reduction in manual processing time for faculty.",
        "Instant Transparency for parents and students.",
        "Data-Driven Insights to help identify students needing extra support."
    ])

    # 8. Conclusion & Future Scope
    add_slide("Conclusion & Future Scope", bullet_points=[
        "Conclusion: A scalable, modern solution for academic data management.",
        "Future Scope: AI-driven performance predictions and integration with LMS systems."
    ])

    prs.save('Student_Assessment_Tracker_Presentation.pptx')
    print("Presentation saved successfully as 'Student_Assessment_Tracker_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
